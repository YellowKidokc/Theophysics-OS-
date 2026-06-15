"""Auto File Sorter — AI-powered classification, naming, and organization.

Tiered intelligence:
  Tier 1: YAKE keywords + rule-based classification (free, instant)
  Tier 2: Ollama local LLM if available (free, slower)
  Tier 3: Claude API if configured (paid, best quality)

Usage:
    python auto_sort.py classify <path>                — classify all files, show proposals
    python auto_sort.py organize <path> <output>       — classify + sort into domain folders
    python auto_sort.py classify-one <filepath>        — classify a single file in detail
"""

import hashlib
import json
import os
import re
import shutil
import sqlite3
import sys
from collections import defaultdict
from datetime import datetime
from pathlib import Path
from typing import Optional

# ─── TEXT EXTRACTION ─────────────────────────────────────────────────────────

def extract_text(filepath: str, max_chars: int = 5000) -> str:
    """Extract text content from a file. Supports txt, md, py, json, csv, pdf, docx."""
    path = Path(filepath)
    ext = path.suffix.lower()

    try:
        # Plain text types
        if ext in {'.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.css',
                   '.html', '.json', '.yaml', '.yml', '.xml', '.csv', '.ini',
                   '.cfg', '.toml', '.log', '.bat', '.sh', '.ps1', '.sql',
                   '.rs', '.go', '.java', '.c', '.cpp', '.h', '.r', '.lean'}:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()[:max_chars]

        # PDF
        if ext == '.pdf':
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(filepath)
                text = ''
                for page in doc:
                    text += page.get_text()
                    if len(text) > max_chars:
                        break
                doc.close()
                return text[:max_chars]
            except ImportError:
                return ''

        # DOCX
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(filepath)
                text = '\n'.join(p.text for p in doc.paragraphs)
                return text[:max_chars]
            except ImportError:
                return ''

        # XLSX
        if ext in {'.xlsx', '.xls'}:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                text = ''
                for ws in wb.worksheets[:3]:
                    for row in ws.iter_rows(max_row=50, values_only=True):
                        text += ' '.join(str(c) for c in row if c) + '\n'
                wb.close()
                return text[:max_chars]
            except ImportError:
                return ''

        # PPTX
        if ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + '\n'
                return text[:max_chars]
            except ImportError:
                return ''

    except Exception:
        return ''

    return ''


# ─── KEYWORD EXTRACTION (TIER 1: YAKE) ───────────────────────────────────────

def extract_keywords_yake(text: str, top_n: int = 10) -> list:
    """Extract keywords using YAKE (no ML, pure statistical)."""
    try:
        import yake
        kw_extractor = yake.KeywordExtractor(
            lan="en", n=2, dedupLim=0.7, top=top_n, features=None
        )
        keywords = kw_extractor.extract_keywords(text)
        return [{'keyword': kw, 'score': round(1 - score, 3), 'source': 'yake'}
                for kw, score in keywords]
    except ImportError:
        # Fallback: simple word frequency
        words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
        stopwords = {'the', 'and', 'for', 'that', 'this', 'with', 'are', 'was',
                     'from', 'have', 'has', 'not', 'but', 'can', 'will', 'been',
                     'which', 'their', 'would', 'there', 'what', 'about', 'when',
                     'make', 'than', 'them', 'into', 'could', 'other', 'more',
                     'also', 'its', 'over', 'such', 'only', 'some', 'very', 'just'}
        freq = defaultdict(int)
        for w in words:
            if w not in stopwords:
                freq[w] += 1
        top = sorted(freq.items(), key=lambda x: -x[1])[:top_n]
        return [{'keyword': w, 'score': round(c / max(1, len(words)), 3), 'source': 'freq'}
                for w, c in top]


# ─── DOMAIN CLASSIFICATION ───────────────────────────────────────────────────

# Domain definitions — keyword signals that map to domains
DOMAIN_RULES = {
    'THEOPHYSICS': {
        'code': 'TP',
        'keywords': ['theophysics', 'master equation', 'coherence', 'decoherence',
                     'chi field', 'logos', 'axiom', 'isomorphism', 'grace operator',
                     'entropy', 'moral conservation', 'noether', 'lagrangian',
                     'trinity', 'christological', 'lowe coherence', 'ten laws',
                     'cross-domain', 'theological', 'scripture', 'genesis',
                     'quantum theology', 'faith', 'prayer', 'gospel', 'biblical',
                     'jesus', 'christ', 'spirit', 'holy', 'church', 'sermon',
                     'worship', 'salvation', 'resurrection', 'prophecy', 'revelation'],
        'extensions': [],
    },
    'DATA_TRADING': {
        'code': 'DT',
        'keywords': ['spy', 'theta', 'options', 'trading', 'stock', 'portfolio',
                     'puts', 'calls', 'strike', 'expiration', 'premium', 'delta',
                     'vix', 'bull', 'bear', 'forex', 'crypto', 'bitcoin', 'ethereum',
                     'candlestick', 'chart', 'technical analysis', 'earnings'],
        'extensions': [],
    },
    'BUSINESS': {
        'code': 'BZ',
        'keywords': ['revenue', 'profit', 'marketing', 'sales', 'customer',
                     'startup', 'business plan', 'roi', 'kpi', 'budget',
                     'invoice', 'contract', 'proposal', 'client', 'vendor',
                     'supply chain', 'ecommerce', 'shopify', 'amazon', 'seo',
                     'affiliate', 'dropship', 'wholesale', 'retail'],
        'extensions': [],
    },
    'DEVELOPMENT': {
        'code': 'DV',
        'keywords': ['function', 'class', 'import', 'return', 'variable',
                     'api', 'endpoint', 'database', 'server', 'docker',
                     'kubernetes', 'git', 'deploy', 'react', 'python',
                     'javascript', 'typescript', 'node', 'npm', 'pip',
                     'cloudflare', 'worker', 'electron', 'fastapi', 'flask',
                     'postgres', 'sqlite', 'mongodb', 'redis'],
        'extensions': {'.py', '.js', '.ts', '.jsx', '.tsx', '.rs', '.go',
                      '.java', '.c', '.cpp', '.h', '.css', '.html', '.sql',
                      '.dockerfile', '.lean'},
    },
    'INFRASTRUCTURE': {
        'code': 'IF',
        'keywords': ['proxmox', 'synology', 'nas', 'homelab', 'bios', 'nvme',
                     'ssd', 'raid', 'backup', 'cloudflare tunnel', 'dns',
                     'firewall', 'network', 'router', 'switch', 'vlan',
                     'pfsense', 'virtualization', 'vm', 'container'],
        'extensions': [],
    },
    'AI_ML': {
        'code': 'AI',
        'keywords': ['model', 'training', 'inference', 'transformer', 'llm',
                     'neural', 'embedding', 'vector', 'prompt', 'fine-tune',
                     'dataset', 'classification', 'nlp', 'gpt', 'claude',
                     'gemini', 'ollama', 'huggingface', 'tensor', 'pytorch',
                     'rag', 'agent', 'langchain'],
        'extensions': {'.gguf', '.safetensors', '.onnx'},
    },
    'MEDIA': {
        'code': 'MD',
        'keywords': ['podcast', 'video', 'audio', 'recording', 'stream',
                     'youtube', 'tts', 'transcript', 'subtitle', 'render',
                     'edit', 'premiere', 'davinci', 'obs', 'thumbnail'],
        'extensions': {'.mp4', '.mkv', '.avi', '.mov', '.mp3', '.wav',
                      '.flac', '.ogg', '.m4a', '.srt', '.vtt'},
    },
    'DOCUMENTS': {
        'code': 'DC',
        'keywords': ['report', 'memo', 'letter', 'template', 'form',
                     'application', 'resume', 'cover letter', 'manual',
                     'guide', 'handbook', 'policy', 'procedure'],
        'extensions': {'.pdf', '.docx', '.doc', '.pptx', '.ppt',
                      '.xlsx', '.xls', '.odt'},
    },
    'IMAGES': {
        'code': 'IM',
        'keywords': ['photo', 'screenshot', 'diagram', 'chart', 'logo',
                     'icon', 'wallpaper', 'banner', 'infographic'],
        'extensions': {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp',
                      '.svg', '.ico', '.tiff', '.psd', '.ai'},
    },
    'PERSONAL': {
        'code': 'PR',
        'keywords': ['journal', 'diary', 'personal', 'family', 'health',
                     'medical', 'insurance', 'tax', 'receipt', 'bill',
                     'bank', 'mortgage', 'rent', 'utility'],
        'extensions': [],
    },
}


def classify_by_keywords(keywords: list, ext: str) -> dict:
    """Rule-based classification using keyword matching against domain rules."""
    scores = {}
    ext = ext.lower()

    for domain, rules in DOMAIN_RULES.items():
        score = 0
        matched = []

        # Extension match (strong signal)
        if rules.get('extensions') and ext in rules['extensions']:
            score += 30
            matched.append(f'ext:{ext}')

        # Keyword matching
        kw_texts = [k['keyword'].lower() for k in keywords]
        full_text = ' '.join(kw_texts)

        for rule_kw in rules['keywords']:
            rule_kw_lower = rule_kw.lower()
            # Exact keyword match
            if rule_kw_lower in kw_texts:
                score += 15
                matched.append(rule_kw)
            # Substring in any keyword
            elif any(rule_kw_lower in kw for kw in kw_texts):
                score += 8
                matched.append(f'~{rule_kw}')
            # Substring in full text
            elif rule_kw_lower in full_text:
                score += 4
                matched.append(f'≈{rule_kw}')

        if score > 0:
            scores[domain] = {
                'score': min(score, 100),
                'code': rules['code'],
                'matched': matched[:5],
            }

    if not scores:
        return {
            'domain': 'UNCATEGORIZED',
            'code': 'UC',
            'confidence': 0,
            'matched': [],
            'all_scores': {},
        }

    best = max(scores.items(), key=lambda x: x[1]['score'])
    return {
        'domain': best[0],
        'code': best[1]['code'],
        'confidence': best[1]['score'],
        'matched': best[1]['matched'],
        'all_scores': {k: v['score'] for k, v in sorted(scores.items(), key=lambda x: -x[1]['score'])},
    }


# ─── LLM CLASSIFICATION (TIER 2/3) ──────────────────────────────────────────

def classify_with_ollama(text: str, filename: str) -> Optional[dict]:
    """Use local Ollama for classification if available."""
    try:
        import urllib.request
        prompt = f"""Classify this file into ONE category. Return ONLY valid JSON.

Categories: THEOPHYSICS, DATA_TRADING, BUSINESS, DEVELOPMENT, INFRASTRUCTURE, AI_ML, MEDIA, DOCUMENTS, IMAGES, PERSONAL

Filename: {filename}
Content preview: {text[:1500]}

Return JSON: {{"domain": "CATEGORY", "confidence": 0-100, "reason": "brief reason", "suggested_name": "descriptive_filename"}}"""

        payload = json.dumps({
            "model": "llama3.2",
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0.2, "num_predict": 200}
        }).encode()

        req = urllib.request.Request(
            "http://localhost:11434/api/generate",
            data=payload,
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read().decode())
            response_text = result.get('response', '')
            # Try to parse JSON from response
            match = re.search(r'\{[^}]+\}', response_text)
            if match:
                return json.loads(match.group())
    except Exception:
        return None
    return None


# ─── INTELLIGENT NAMING ──────────────────────────────────────────────────────

def generate_slug(keywords: list, max_len: int = 40) -> str:
    """Generate a descriptive slug from top keywords."""
    if not keywords:
        return 'unnamed'
    top = [k['keyword'] for k in keywords[:3]]
    slug = '-'.join(top)
    slug = re.sub(r'[^a-zA-Z0-9\-]', '', slug.lower())
    slug = re.sub(r'-+', '-', slug).strip('-')
    return slug[:max_len] or 'unnamed'


def generate_proposed_name(filename: str, classification: dict, keywords: list, seq_id: int) -> str:
    """Generate FIS-style proposed filename."""
    ext = Path(filename).suffix.lower()
    slug = generate_slug(keywords)
    code = classification['code']
    return f"{slug}_{code}_{seq_id:06d}{ext}"


# ─── SQLITE STORAGE ──────────────────────────────────────────────────────────

DB_PATH = None

def get_db(root: str) -> sqlite3.Connection:
    """Get or create SQLite database for this sort session."""
    global DB_PATH
    DB_PATH = os.path.join(root, '_autosort.db')
    conn = sqlite3.connect(DB_PATH)
    conn.execute('''CREATE TABLE IF NOT EXISTS classified_files (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        original_name TEXT,
        original_path TEXT,
        proposed_name TEXT,
        domain TEXT,
        domain_code TEXT,
        confidence REAL,
        keywords TEXT,
        matched_rules TEXT,
        sha256 TEXT,
        file_size INTEGER,
        status TEXT DEFAULT 'pending',
        created_at TEXT DEFAULT CURRENT_TIMESTAMP
    )''')
    conn.commit()
    return conn


def md5_quick(filepath: str) -> str:
    """Quick hash for dedup."""
    h = hashlib.md5()
    with open(filepath, 'rb') as f:
        h.update(f.read(65536))  # First 64KB only for speed
    return h.hexdigest()


# ─── MAIN PIPELINE ───────────────────────────────────────────────────────────

def classify_file(filepath: str, use_llm: bool = False, use_nlp: bool = False,
                   use_markov: bool = True) -> dict:
    """Classify a single file through the tiered pipeline.
    
    Tier 1: YAKE keywords + rule matching (instant, always runs)
    Tier 2: DeBERTa NLI + BART summarizer (--nlp flag, fires on low confidence)
    Tier 3: Markov preference engine (always runs if trained, predicts your vote)
    Tier 4: Ollama local LLM (--llm flag, last resort)
    """
    path = Path(filepath)
    if not path.exists():
        return {'error': f'File not found: {filepath}'}

    # Extract text
    text = extract_text(filepath)
    filename = path.name
    ext = path.suffix.lower()
    size = path.stat().st_size

    # Tier 1: YAKE + rules (always runs)
    keywords = extract_keywords_yake(text) if text.strip() else []
    classification = classify_by_keywords(keywords, ext)
    classification['source'] = 'yake'

    # Tier 2: NLP models (DeBERTa + BART) for low confidence
    nlp_result = None
    nlp_summary = None
    if use_nlp and classification['confidence'] < 50 and text.strip():
        try:
            from nlp_bridge import classify_with_deberta, summarize_with_bart, classify_image_with_clip
            
            # DeBERTa zero-shot classification
            deberta = classify_with_deberta(text)
            if deberta and deberta['confidence'] > classification['confidence']:
                classification['domain'] = deberta['domain']
                classification['code'] = DOMAIN_RULES.get(deberta['domain'], {}).get('code', 'UC')
                classification['confidence'] = deberta['confidence']
                classification['matched'] = [f"deberta:{deberta['label']}"]
                classification['source'] = 'deberta'
                nlp_result = deberta
            
            # BART summary for naming
            summary = summarize_with_bart(text)
            if summary:
                nlp_summary = summary
            
            # CLIP for images
            if ext in {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}:
                clip = classify_image_with_clip(filepath)
                if clip and clip['confidence'] > classification['confidence']:
                    nlp_result = clip
                    
        except ImportError:
            pass  # nlp_bridge not available
        except Exception as e:
            print(f"  [NLP] Error: {e}")

    # Tier 3: Markov preference engine (learned from your decisions)
    markov_prediction = None
    if use_markov:
        try:
            from preference_engine import predict_domain
            kw_list = [k['keyword'] for k in keywords]
            markov_prediction = predict_domain(kw_list, ext)
            
            # If Markov is more confident AND disagrees, it wins
            if (markov_prediction and 
                markov_prediction['confidence'] > classification['confidence'] and
                markov_prediction['domain'] != classification['domain']):
                classification['_original_domain'] = classification['domain']
                classification['_original_confidence'] = classification['confidence']
                classification['domain'] = markov_prediction['domain']
                classification['code'] = DOMAIN_RULES.get(markov_prediction['domain'], {}).get('code', 'UC')
                classification['confidence'] = markov_prediction['confidence']
                classification['source'] = 'markov'
                classification['matched'] = [f"markov:learned_preference"]
        except ImportError:
            pass
        except Exception as e:
            print(f"  [Markov] Error: {e}")

    # Tier 4: Ollama LLM (last resort for still-low confidence)
    llm_result = None
    if use_llm and classification['confidence'] < 50:
        llm_result = classify_with_ollama(text, filename)
        if llm_result and llm_result.get('confidence', 0) > classification['confidence']:
            domain_upper = llm_result.get('domain', '').upper()
            if domain_upper in DOMAIN_RULES:
                classification['domain'] = domain_upper
                classification['code'] = DOMAIN_RULES[domain_upper]['code']
                classification['confidence'] = llm_result['confidence']
                classification['matched'] = [f"llm:{llm_result.get('reason', '')}"]
                classification['source'] = 'ollama'

    return {
        'filepath': str(path),
        'filename': filename,
        'ext': ext,
        'size': size,
        'text_preview': text[:200] if text else '',
        'keywords': keywords,
        'classification': classification,
        'llm_result': llm_result,
        'nlp_result': nlp_result,
        'nlp_summary': nlp_summary,
        'markov_prediction': markov_prediction,
    }


IGNORE_DIRS = {
    'node_modules', '__pycache__', '.git', '.svn', '.hg', '.vs', '.vscode',
    'deps', 'dist', 'build', 'out', '.next', '.nuxt', '.cache', '.tmp',
    'vendor', 'target', 'bin', 'obj', '.idea', '.gradle', 'venv', '.venv',
    'env', '.env', '.tox', 'egg-info', '.eggs', '.mypy_cache', '.pytest_cache',
    'coverage', '.nyc_output', 'bower_components', '.sass-cache',
    '$Recycle.Bin', 'System Volume Information', '.Trash',
}

IGNORE_FILES = {
    '_autosort.db', 'preference_engine.db', '.sortconfig.yaml',
    'thumbs.db', 'desktop.ini', '.ds_store', 'icon\r',
}


def classify_directory(root: str, use_llm: bool = False, use_nlp: bool = False,
                       use_markov: bool = True, extensions: set = None,
                       top_level_only: bool = False) -> list:
    """Classify all files in a directory."""
    root_path = Path(root)
    results = []
    files = []

    if top_level_only:
        for item in root_path.iterdir():
            if item.is_file() and not item.name.startswith('.') and item.name.lower() not in IGNORE_FILES:
                if extensions and item.suffix.lower() not in extensions:
                    continue
                files.append(str(item))
    else:
        for dirpath, dirnames, filenames in os.walk(root_path):
            # Skip ignored directories
            dirnames[:] = [d for d in dirnames if d not in IGNORE_DIRS and not d.startswith('.')]
            for fname in filenames:
                if fname.startswith('.') or fname.lower() in IGNORE_FILES:
                    continue
                fpath = Path(dirpath) / fname
                if extensions and fpath.suffix.lower() not in extensions:
                    continue
                files.append(str(fpath))

    print(f"\n  Classifying {len(files)} files...")
    for i, fpath in enumerate(files):
        if (i + 1) % 25 == 0 or i == 0:
            print(f"    [{i+1}/{len(files)}] {Path(fpath).name[:50]}")
        result = classify_file(fpath, use_llm=use_llm, use_nlp=use_nlp, use_markov=use_markov)
        if 'error' not in result:
            results.append(result)

    return results


def print_classification(results: list):
    """Pretty-print classification results."""
    # Group by domain
    by_domain = defaultdict(list)
    for r in results:
        domain = r['classification']['domain']
        by_domain[domain].append(r)

    print(f"\n{'='*70}")
    print(f"  AUTO-CLASSIFICATION RESULTS: {len(results)} files")
    print(f"{'='*70}")

    for domain in sorted(by_domain.keys()):
        files = by_domain[domain]
        print(f"\n  [{domain}] — {len(files)} files")
        for f in files[:10]:
            conf = f['classification']['confidence']
            bar = '#' * int(conf / 5) + '.' * (20 - int(conf / 5))
            kws = ', '.join(k['keyword'] for k in f['keywords'][:3])
            print(f"    {bar} {conf:3.0f}%  {f['filename'][:40]:<40s}  [{kws}]")
        if len(files) > 10:
            print(f"    ... +{len(files)-10} more")

    # Confidence distribution
    confs = [r['classification']['confidence'] for r in results]
    high = sum(1 for c in confs if c >= 70)
    mid = sum(1 for c in confs if 30 <= c < 70)
    low = sum(1 for c in confs if c < 30)
    print(f"\n  Confidence: {high} high (>=70) | {mid} medium (30-70) | {low} low (<30)")
    print(f"{'='*70}\n")


def organize_files(root: str, output: str, use_llm: bool = False,
                   dry_run: bool = False, extensions: set = None):
    """Classify files and sort them into domain folders."""
    results = classify_directory(root, use_llm=use_llm, extensions=extensions)
    if not results:
        print("  No files to organize.")
        return

    print_classification(results)

    # Build operations
    output_path = Path(output)
    ops = []
    seq = 1

    for r in results:
        domain = r['classification']['domain']
        proposed = generate_proposed_name(r['filename'], r['classification'], r['keywords'], seq)
        dest_dir = output_path / domain
        dest_file = dest_dir / proposed

        ops.append({
            'src': r['filepath'],
            'dest': str(dest_file),
            'dir': str(dest_dir),
            'domain': domain,
            'confidence': r['classification']['confidence'],
            'proposed_name': proposed,
            'original_name': r['filename'],
        })
        seq += 1

    # Summary
    dirs_needed = set(o['dir'] for o in ops)
    print(f"  Organization plan: {len(ops)} files → {len(dirs_needed)} domain folders")
    for d in sorted(dirs_needed):
        count = sum(1 for o in ops if o['dir'] == d)
        print(f"    {Path(d).name:20s} → {count:>5} files")

    if dry_run:
        print("\n  [DRY RUN — no files moved]")
        # Save to SQLite anyway for review
        conn = get_db(root)
        for o in ops:
            r = next(x for x in results if x['filepath'] == o['src'])
            conn.execute(
                '''INSERT INTO classified_files
                   (original_name, original_path, proposed_name, domain, domain_code,
                    confidence, keywords, matched_rules, file_size, status)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)''',
                (o['original_name'], o['src'], o['proposed_name'], o['domain'],
                 r['classification']['code'], o['confidence'],
                 json.dumps([k['keyword'] for k in r['keywords'][:5]]),
                 json.dumps(r['classification'].get('matched', [])),
                 r['size'], 'proposed')
            )
        conn.commit()
        conn.close()
        print(f"  Proposals saved to {DB_PATH}")
        return ops

    confirm = input(f"\n  Sort {len(ops)} files into domain folders? (y/n): ").strip().lower()
    if confirm != 'y':
        print("  Cancelled.")
        return []

    for o in ops:
        Path(o['dir']).mkdir(parents=True, exist_ok=True)
        shutil.copy2(o['src'], o['dest'])

    print(f"  Done. {len(ops)} files organized into {len(dirs_needed)} domains.")
    return ops


# ─── CLI ─────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1].lower()
    
    # Stats command — no path needed
    if cmd == 'stats':
        try:
            from preference_engine import get_engine_stats
            stats = get_engine_stats()
            print(f"\n{'='*50}")
            print(f"  PREFERENCE ENGINE STATS")
            print(f"{'='*50}")
            print(f"  Total decisions:     {stats['total_decisions']}")
            print(f"  Approvals:           {stats['approvals']}")
            print(f"  Overrides:           {stats['overrides']}")
            print(f"  Rejects:             {stats['rejects']}")
            print(f"  Accuracy:            {stats['accuracy']}%")
            print(f"  Keywords learned:    {stats['unique_keywords_learned']}")
            print(f"  Auto-approve at:     {stats['auto_approve_threshold']}%")
            if stats['top_corrections']:
                print(f"\n  Top corrections:")
                for c in stats['top_corrections']:
                    print(f"    {c['from']:20s} → {c['to']:20s}  ({c['count']}x)")
            print(f"{'='*50}\n")
        except Exception as e:
            print(f"  Error: {e}")
        return

    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    path = sys.argv[2]
    use_llm = '--llm' in sys.argv
    use_nlp = '--nlp' in sys.argv
    use_markov = '--no-markov' not in sys.argv
    dry_run = '--dry' in sys.argv
    top_only = '--top' in sys.argv

    # Parse --ext filter
    extensions = None
    if '--ext' in sys.argv:
        idx = sys.argv.index('--ext')
        if idx + 1 < len(sys.argv):
            extensions = set('.' + e.lstrip('.') for e in sys.argv[idx + 1].split(','))

    if not os.path.exists(path):
        print(f"  Error: path does not exist: {path}")
        sys.exit(1)

    if cmd == 'classify':
        results = classify_directory(path, use_llm=use_llm, use_nlp=use_nlp,
                                     use_markov=use_markov, extensions=extensions,
                                     top_level_only=top_only)
        print_classification(results)

    elif cmd == 'classify-one':
        result = classify_file(path, use_llm=use_llm, use_nlp=use_nlp, use_markov=use_markov)
        if 'error' in result:
            print(f"  Error: {result['error']}")
        else:
            print(f"\n  File: {result['filename']}")
            print(f"  Size: {result['size']:,} bytes")
            print(f"  Domain: {result['classification']['domain']} ({result['classification']['code']})")
            print(f"  Confidence: {result['classification']['confidence']}%")
            print(f"  Source: {result['classification'].get('source', 'yake')}")
            print(f"  Matched: {result['classification']['matched']}")
            print(f"  Keywords: {[k['keyword'] for k in result['keywords'][:8]]}")
            print(f"  All scores: {result['classification']['all_scores']}")
            if result.get('markov_prediction'):
                mp = result['markov_prediction']
                print(f"  Markov says: {mp['domain']} ({mp['confidence']}%, trained on {mp['training_size']} decisions)")
                if mp.get('correction_warning'):
                    cw = mp['correction_warning']
                    print(f"  ⚠ Override warning: {cw['override_rate']}% corrected to {cw['usually_corrected_to']}")
            if result.get('nlp_result'):
                print(f"  NLP result: {result['nlp_result']}")
            if result.get('nlp_summary'):
                print(f"  NLP summary: {result['nlp_summary']}")
            if result.get('text_preview'):
                print(f"  Preview: {result['text_preview'][:150]}...")

    elif cmd == 'learn':
        # Manual learning: python auto_sort.py learn <file> <domain>
        if len(sys.argv) < 4:
            print("  Usage: auto_sort.py learn <filepath> <DOMAIN>")
            sys.exit(1)
        domain = sys.argv[3].upper()
        result = classify_file(path, use_llm=False, use_nlp=False, use_markov=False)
        if 'error' in result:
            print(f"  Error: {result['error']}")
            return
        try:
            from preference_engine import record_decision
            keywords = [k['keyword'] for k in result['keywords']]
            proposed = result['classification']['domain']
            action = 'approve' if proposed == domain else 'override'
            record_decision(
                filename=result['filename'],
                extension=result['ext'],
                keywords=keywords,
                proposed_domain=proposed,
                final_domain=domain,
                confidence=result['classification']['confidence'],
                action=action,
                source=result['classification'].get('source', 'yake'),
            )
            print(f"  Learned: {result['filename']} → {domain} ({action})")
        except Exception as e:
            print(f"  Error recording: {e}")

    elif cmd == 'organize':
        output = sys.argv[3] if len(sys.argv) > 3 and not sys.argv[3].startswith('--') else os.path.join(path, '_organized')
        organize_files(path, output, use_llm=use_llm, dry_run=dry_run, extensions=extensions)

    elif cmd == 'rename':
        rename_in_place(path, use_llm=use_llm, dry_run=dry_run, extensions=extensions)

    elif cmd == 'folders':
        results = classify_folders(path, use_llm=use_llm)
        print_folder_classification(results)

    elif cmd == 'organize-folders':
        output = sys.argv[3] if len(sys.argv) > 3 and not sys.argv[3].startswith('--') else os.path.join(path, '_organized')
        organize_folders(path, output, use_llm=use_llm, dry_run=dry_run)

    else:
        print(f"  Unknown command: {cmd}")
        print(__doc__)
        sys.exit(1)


def rename_in_place(root: str, use_llm: bool = False, dry_run: bool = True,
                    extensions: set = None):
    """Classify files and rename them using FIS naming schema WITHOUT moving them.
    Files stay in their current folders — only the names change."""
    results = classify_directory(root, use_llm=use_llm, extensions=extensions)
    if not results:
        print("  No files to rename.")
        return

    print_classification(results)

    # Build rename operations
    ops = []
    seq = 1
    for r in results:
        path = Path(r['filepath'])
        proposed = generate_proposed_name(r['filename'], r['classification'], r['keywords'], seq)
        dest = path.parent / proposed
        if str(dest) != str(path):
            ops.append({
                'src': str(path),
                'dest': str(dest),
                'original': r['filename'],
                'proposed': proposed,
                'domain': r['classification']['domain'],
                'confidence': r['classification']['confidence'],
            })
        seq += 1

    print(f"\n  Rename proposals: {len(ops)} files")
    for o in ops[:20]:
        print(f"    {o['original'][:35]:<35s} -> {o['proposed']}")
    if len(ops) > 20:
        print(f"    ... +{len(ops) - 20} more")

    if dry_run:
        print("\n  [DRY RUN - no files renamed]")
        return ops

    confirm = input(f"\n  Rename {len(ops)} files in place? (y/n): ").strip().lower()
    if confirm != 'y':
        print("  Cancelled.")
        return []

    renamed = 0
    for o in ops:
        try:
            Path(o['src']).rename(o['dest'])
            renamed += 1
        except OSError as e:
            print(f"    Error: {o['original']}: {e}")

    print(f"  Done. {renamed} files renamed.")
    return ops


# ─── FOLDER-LEVEL CLASSIFICATION ─────────────────────────────────────────────

def classify_folders(root: str, use_llm: bool = False) -> list:
    """Classify top-level folders as units based on their contents."""
    root_path = Path(root)
    folder_results = []

    folders = [f for f in root_path.iterdir() if f.is_dir() and not f.name.startswith('.')]
    print(f"\n  Classifying {len(folders)} folders...")

    for folder in sorted(folders):
        # Sample files from folder (up to 20 for speed)
        sample_files = []
        for dirpath, _, filenames in os.walk(folder):
            for fname in filenames:
                fpath = Path(dirpath) / fname
                if not fname.startswith('.') and fpath.suffix.lower() not in {'.pyc', '.o', '.obj'}:
                    sample_files.append(str(fpath))
                    if len(sample_files) >= 20:
                        break
            if len(sample_files) >= 20:
                break

        # Classify each sample file
        file_classifications = []
        for sf in sample_files:
            r = classify_file(sf, use_llm=use_llm)
            if 'error' not in r:
                file_classifications.append(r)

        # Aggregate: majority vote on domain
        if file_classifications:
            domain_votes = defaultdict(float)
            all_keywords = []
            for fc in file_classifications:
                d = fc['classification']['domain']
                c = fc['classification']['confidence']
                domain_votes[d] += c
                all_keywords.extend(fc['keywords'][:3])

            best_domain = max(domain_votes.items(), key=lambda x: x[1])
            avg_conf = best_domain[1] / len(file_classifications)
            domain_code = DOMAIN_RULES.get(best_domain[0], {}).get('code', 'UC')
        else:
            best_domain = ('UNCATEGORIZED', 0)
            avg_conf = 0
            all_keywords = []
            domain_code = 'UC'

        # Count files
        total = sum(1 for _ in folder.rglob('*') if _.is_file())

        folder_results.append({
            'folder': str(folder),
            'name': folder.name,
            'domain': best_domain[0],
            'domain_code': domain_code,
            'confidence': round(avg_conf, 1),
            'total_files': total,
            'sample_size': len(file_classifications),
            'keywords': list(set(k['keyword'] for k in all_keywords[:10])),
        })

    return folder_results


def print_folder_classification(results: list):
    """Pretty-print folder classification."""
    print(f"\n{'='*70}")
    print(f"  FOLDER CLASSIFICATION: {len(results)} folders")
    print(f"{'='*70}")
    for r in sorted(results, key=lambda x: x['domain']):
        conf = r['confidence']
        bar = '#' * int(conf / 5) + '.' * (20 - int(conf / 5))
        kws = ', '.join(r['keywords'][:3])
        print(f"    {bar} {conf:3.0f}%  [{r['domain_code']}] {r['name'][:30]:<30s}  ({r['total_files']} files)  [{kws}]")
    print(f"{'='*70}\n")


def organize_folders(root: str, output: str, use_llm: bool = False, dry_run: bool = True):
    """Classify folders and move them as units into domain parent folders."""
    results = classify_folders(root, use_llm=use_llm)
    if not results:
        print("  No folders to organize.")
        return

    print_folder_classification(results)

    output_path = Path(output)
    ops = []
    for r in results:
        dest = output_path / r['domain'] / r['name']
        if str(dest) != r['folder']:
            ops.append({
                'src': r['folder'],
                'dest': str(dest),
                'domain': r['domain'],
                'files': r['total_files'],
            })

    print(f"\n  Move plan: {len(ops)} folders")
    for o in ops:
        print(f"    {Path(o['src']).name:<30s} -> {o['domain']}/  ({o['files']} files)")

    if dry_run:
        print("\n  [DRY RUN - no folders moved]")
        return ops

    confirm = input(f"\n  Move {len(ops)} folders? (y/n): ").strip().lower()
    if confirm != 'y':
        print("  Cancelled.")
        return []

    for o in ops:
        Path(o['dest']).parent.mkdir(parents=True, exist_ok=True)
        shutil.move(o['src'], o['dest'])

    print(f"  Done. {len(ops)} folders organized.")
    return ops



# ─── ENTRY POINT ─────────────────────────────────────────────────────────────

if __name__ == '__main__':
    main()
